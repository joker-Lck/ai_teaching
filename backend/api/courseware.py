"""
课件生成 API - WebSocket 实时进度 + PPTX/DOCX 下载
"""
import json
import io
import os
import time
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse, Response
from typing import Optional

from backend.schemas.models import (
    CoursewareGenerateRequest, CoursewareResponse,
    ClarificationRequest, ClarificationReplyRequest,
    RefineRequest, BaseResponse,
)
from backend.dependencies import get_current_user, get_api_config
from services.courseware_service import CoursewareService
from data.db_operations import db
from core.logger import info, error as log_error

router = APIRouter()


@router.post("/generate", response_model=CoursewareResponse)
async def generate_courseware(
    req: CoursewareGenerateRequest,
    user: dict = Depends(get_current_user),
):
    """生成课件

    同步接口, 返回完整的课件数据。
    前端可通过 WebSocket /ws/generation 获取实时进度。
    """
    config = get_api_config()

    courseware_service = CoursewareService()

    try:
        result = courseware_service.generate_courseware(
            topic=req.topic,
            requirements_text=req.requirements if req.requirements else "无特殊要求",
            api_key=config["api_key"],
            base_url=config["base_url"],
            fast_mode=req.fast_mode,
        )

        if result["success"]:
            return CoursewareResponse(
                success=True,
                subject=result["subject"],
                outline=result["outline"],
                slides=result["slides"],
                theme=result["theme"],
                courseware_id=result.get("courseware_id"),
                generated_images=result.get("generated_images", {}),
            )
        else:
            return CoursewareResponse(
                success=False,
                error=result.get("error", "生成失败"),
            )

    except Exception as e:
        log_error(f"课件生成异常: {str(e)}")
        return CoursewareResponse(success=False, error=str(e))


@router.post("/generate-stream")
async def generate_courseware_stream(
    req: CoursewareGenerateRequest,
    user: dict = Depends(get_current_user),
):
    """生成课件 - SSE 流式进度推送"""
    config = get_api_config()

    async def generate_progress():
        courseware_service = CoursewareService()

        try:
            # 步骤 1: 学科识别
            yield f"data: {json.dumps({'step': 1, 'total': 3, 'message': '正在识别学科并生成大纲...', 'progress': 10}, ensure_ascii=False)}\n\n"

            result = courseware_service.generate_courseware(
                topic=req.topic,
                requirements_text=req.requirements if req.requirements else "无特殊要求",
                api_key=config["api_key"],
                base_url=config["base_url"],
                fast_mode=req.fast_mode,
            )

            if result["success"]:
                yield f"data: {json.dumps({'step': 2, 'total': 3, 'message': 'PPT 结构生成完成', 'progress': 70}, ensure_ascii=False)}\n\n"
                yield f"data: {json.dumps({'step': 3, 'total': 3, 'message': '保存到数据库', 'progress': 90}, ensure_ascii=False)}\n\n"

                # 发送最终结果
                yield f"data: {json.dumps({\
                    'type': 'done',\
                    'subject': result['subject'],\
                    'outline': result['outline'],\
                    'slides': result['slides'],\
                    'theme': result['theme'],\
                    'courseware_id': result.get('courseware_id'),\
                    'progress': 100,\
                }, ensure_ascii=False)}\n\n"
            else:
                yield f"data: {json.dumps({'type': 'error', 'message': result.get('error', '生成失败')}, ensure_ascii=False)}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(
        generate_progress(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@router.post("/identify")
async def identify_subject(
    req: CoursewareGenerateRequest,
    user: dict = Depends(get_current_user),
):
    """学科识别"""
    config = get_api_config()
    courseware_service = CoursewareService()

    try:
        from openai import OpenAI
        from core.prompts import CoursewarePrompts
        from core.utils import clean_json_string, safe_json_loads

        client = OpenAI(api_key=config["api_key"], base_url=config["base_url"])
        prompt = CoursewarePrompts.get_identify_prompt(req.topic, req.requirements or "无特殊要求")

        response = client.chat.completions.create(
            model="moonshot-v1-8k",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=1000,
            timeout=30,
        )

        content = clean_json_string(response.choices[0].message.content)
        data = safe_json_loads(content)

        if data:
            return {"success": True, "subject": data.get("subject", "综合"), "outline": data.get("outline", "")}
        else:
            return {"success": False, "error": "学科识别失败"}

    except Exception as e:
        return {"success": False, "error": str(e)}


@router.post("/clarify")
async def clarify_requirements(
    req: ClarificationRequest,
    user: dict = Depends(get_current_user),
):
    """开始需求澄清对话"""
    config = get_api_config()
    courseware_service = CoursewareService()

    result = courseware_service.start_clarification(
        topic=req.topic,
        api_key=config["api_key"],
        base_url=config["base_url"],
    )

    return result


@router.post("/clarify-reply")
async def clarify_reply(
    req: ClarificationReplyRequest,
    user: dict = Depends(get_current_user),
):
    """继续需求澄清对话"""
    config = get_api_config()
    courseware_service = CoursewareService()

    result = courseware_service.continue_clarification(
        topic=req.topic,
        conversation_history=req.conversation_history,
        api_key=config["api_key"],
        base_url=config["base_url"],
    )

    return result


@router.post("/refine", response_model=CoursewareResponse)
async def refine_courseware(
    req: RefineRequest,
    user: dict = Depends(get_current_user),
):
    """基于反馈调整课件"""
    config = get_api_config()
    courseware_service = CoursewareService()

    result = courseware_service.refine_courseware(
        feedback=req.feedback,
        topic=req.topic,
        subject=req.subject,
        slides=req.slides,
        api_key=config["api_key"],
        base_url=config["base_url"],
    )

    if result["success"]:
        return CoursewareResponse(
            success=True,
            slides=result["slides"],
            theme=result.get("theme", {}),
        )
    else:
        return CoursewareResponse(success=False, error=result.get("error", "调整失败"))


@router.get("/history")
async def get_courseware_history(
    user: dict = Depends(get_current_user),
):
    """获取历史课件列表"""
    try:
        courseware_list = db.get_all_courseware()
        items = []
        for cw in (courseware_list or []):
            items.append({
                "id": cw.get("id"),
                "title": cw.get("title", ""),
                "subject": cw.get("subject", ""),
                "grade_level": cw.get("grade_level", ""),
                "created_at": str(cw.get("created_at", "")),
            })
        return {"success": True, "items": items, "total": len(items)}
    except Exception as e:
        return {"success": False, "items": [], "error": str(e)}


@router.get("/history/{courseware_id}")
async def get_courseware_detail(
    courseware_id: int,
    user: dict = Depends(get_current_user),
):
    """获取单个课件详情"""
    try:
        cw = db.get_courseware_by_id(courseware_id)
        if cw:
            content = json.loads(cw.get("content", "{}")) if isinstance(cw.get("content"), str) else cw.get("content", {})
            return {
                "success": True,
                "id": cw.get("id"),
                "title": cw.get("title", ""),
                "subject": cw.get("subject", ""),
                "slides": content.get("slides", []),
                "theme": content.get("theme", {}),
            }
        return {"success": False, "error": "课件不存在"}
    except Exception as e:
        return {"success": False, "error": str(e)}


@router.get("/download/{courseware_id}")
async def download_courseware(
    courseware_id: int,
    format: str = "pptx",
    user: dict = Depends(get_current_user),
):
    """下载课件文件 (PPTX 或 DOCX)"""
    try:
        cw = db.get_courseware_by_id(courseware_id)
        if not cw:
            raise HTTPException(status_code=404, detail="课件不存在")

        content = json.loads(cw.get("content", "{}")) if isinstance(cw.get("content"), str) else cw.get("content", {})
        slides = content.get("slides", [])
        theme = content.get("theme", {})
        title = cw.get("title", "课件")

        if format == "docx":
            # 生成 Word 教案
            from docx import Document
            from docx.shared import Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            heading = doc.add_heading(f"{title} - 教案", 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph(f"学科：{cw.get('subject', '')}")

            for i, slide in enumerate(slides, 1):
                doc.add_heading(f"第{i}页：{slide.get('title', '无标题')}", level=2)
                for point in slide.get("content", []):
                    if point and point.strip():
                        doc.add_paragraph(point, style="List Bullet")

            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)

            return Response(
                content=buf.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f'attachment; filename="{title}_教案.docx"'},
            )
        else:
            # 生成 PPTX
            from pptx import Presentation
            from pptx.util import Pt, Inches
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            primary_color = theme.get("primary_color", "#0a192f")
            accent_color = theme.get("accent_color", "#00d4ff")

            def hex_to_rgb(h):
                h = h.lstrip("#")
                return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))

            for i, slide_data in enumerate(slides):
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 标题栏
                header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
                header.fill.solid()
                header.fill.fore_color.rgb = RGBColor(*hex_to_rgb(primary_color))
                header.line.fill.background()

                # 标题文字
                title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
                p = title_box.text_frame.paragraphs[0]
                p.text = slide_data.get("title", "无标题")
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                # 内容
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
                tf = content_box.text_frame
                tf.word_wrap = True
                for j, point in enumerate(slide_data.get("content", [])):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(10)

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)

            return Response(
                content=buf.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="{title}_课件.pptx"'},
            )

    except HTTPException:
        raise
    except Exception as e:
        log_error(f"下载课件失败: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
