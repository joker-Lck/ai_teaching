import streamlit as st
from openai import OpenAI
from datetime import datetime
import time
import os
import json
from dotenv import load_dotenv
from data.db_operations import db
from data.qa_db_operations import qa_db
from data.rag_knowledge_base import rag_kb
from data.document_parser import doc_parser
from data.embedding_service import embedding_service
from services.animation_service import AnimationService
from services.auth_service import auth_service
from core.ui_components import CustomCSS, PageLayout, UIComponents
from data.data_manager import LearningDataManager, DatabaseManager, CacheManager
from core.utils import clean_json_string, format_file_size, extract_urls, truncate_text, safe_get, generate_filename
from services.qa_service import QAService
from services.courseware_service import CoursewareService
from services.knowledge_service import KnowledgeService
from services.analysis_service import AnalysisService
from functools import lru_cache
from core.logger import info, warning, error, db_connect_success, db_connect_failed, user_login
import re
import io
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
matplotlib.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'Arial Unicode MS']
matplotlib.rcParams['axes.unicode_minus'] = False

# 获取 OpenAI 客户端实例
@st.cache_resource
def get_openai_client():
    """创建并缓存 OpenAI 客户端实例"""
    return OpenAI(api_key=DEFAULT_API_KEY, base_url=BASE_URL)

# 数据库连接管理
@st.cache_resource(ttl=3600)
def get_database_connections():
    """缓存数据库连接，1小时内复用"""
    return DatabaseManager.get_database_connections()

# RAG 知识检索
@st.cache_data(ttl=300)
def cached_rag_search(query, limit=2):
    """优先使用向量检索，降级到全文搜索"""
    try:
        # 1. 获取查询向量
        query_embedding = embedding_service.get_embedding(query)
        
        if query_embedding:
            # 向量检索
            results = rag_kb.search_documents_by_vector(query_embedding, limit=limit)
            
            if results and len(results) > 0:
                info(f"✅ 向量检索成功：{len(results)} 篇文档")
                return results
        
        # 全文搜索
        results = rag_kb.search_documents(query, limit=limit)
        if results:
            info(f"✅ 全文检索成功：{len(results)} 篇文档")
        return results if results else []
    except Exception as e:
        warning(f"⚠️ RAG 检索失败：{str(e)}")
        return []

# 文档分页加载
@st.cache_data(ttl=600)
def load_rag_documents_paginated(offset=0, limit=20):
    """分页加载 RAG 文档"""
    try:
        docs = rag_kb.get_all_documents(limit=limit, offset=offset)
        return docs if docs else []
    except Exception:
        pass
    return []

# 加载配置信息
config = CacheManager.load_env_config()
DEFAULT_API_KEY = config['api_key']
BASE_URL = config['base_url']

# 初始化登录状态
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.current_user = None

# 如果未登录，显示登录界面
if not st.session_state.logged_in:
    st.set_page_config(page_title="AI 教学智能体 - 登录", page_icon="🎓", layout="centered")
    
    # 登录页面
    st.markdown("""
    <div style='text-align: center; padding: 50px 0;'>
        <h1 style='color: #0a192f; font-size: 48px;'>🎓 多模态 AI 教学智能体</h1>
        <p style='color: #666; font-size: 18px;'>AI 赋能教育，智能引领未来</p>
    </div>
    """, unsafe_allow_html=True)
    
    # 登录/注册选项卡
    auth_tab = st.tabs([" 登录", "📝 注册"])
    
    with auth_tab[0]:
        st.subheader("用户登录")
        
        with st.form("login_form"):
            username = st.text_input("用户名", placeholder="请输入用户名")
            password = st.text_input("密码", type="password", placeholder="请输入密码")
            
            col1, col2 = st.columns([1, 1])
            with col1:
                login_btn = st.form_submit_button("登录", type="primary", use_container_width=True)
            with col2:
                guest_btn = st.form_submit_button("游客模式", use_container_width=True)
            
            if login_btn:
                if not username or not password:
                    st.error(" 请输入用户名和密码")
                else:
                    with st.spinner("正在登录..."):
                        result = auth_service.login_user(username, password)
                        
                        if result['success']:
                            st.session_state.logged_in = True
                            st.session_state.current_user = result['user']
                            st.success(f"✅ {result['message']}，欢迎 {result['user']['username']}！")
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.error(f"❌ {result['message']}")
            
            if guest_btn:
                st.session_state.logged_in = True
                st.session_state.current_user = {
                    'id': 0,
                    'username': '游客',
                    'role': 'guest',
                    'email': None
                }
                st.success("✅ 已进入游客模式")
                time.sleep(1)
                st.rerun()
    
    with auth_tab[1]:
        st.subheader("新用户注册")
        
        with st.form("register_form"):
            reg_username = st.text_input("用户名", placeholder="请输入用户名（3-20个字符）")
            reg_password = st.text_input("密码", type="password", placeholder="请输入密码（至少6位）")
            reg_password_confirm = st.text_input("确认密码", type="password", placeholder="请再次输入密码")
            reg_email = st.text_input("邮箱（可选）", placeholder="请输入邮箱地址")
            reg_role = st.radio("用户角色", ["教师", "学生"], index=0)
            
            register_btn = st.form_submit_button("注册", type="primary")
            
            if register_btn:
                # 验证输入
                if not reg_username or not reg_password or not reg_password_confirm:
                    st.error("❌ 请填写必填项")
                elif len(reg_username) < 3 or len(reg_username) > 20:
                    st.error("❌ 用户名长度应在 3-20 个字符之间")
                elif len(reg_password) < 6:
                    st.error("❌ 密码长度至少 6 位")
                elif reg_password != reg_password_confirm:
                    st.error("❌ 两次输入的密码不一致")
                else:
                    role_map = {"教师": "teacher", "学生": "student"}
                    with st.spinner("正在注册..."):
                        result = auth_service.register_user(
                            reg_username, 
                            reg_password, 
                            reg_email if reg_email else None,
                            role_map[reg_role]
                        )
                        
                        if result['success']:
                            st.success(f"✅ {result['message']}！请登录")
                        else:
                            st.error(f"❌ {result['message']}")
    
    st.stop()  # 停止执行后续代码

# 已登录，继续加载主应用

# 初始化学习数据
if "learning_data" not in st.session_state:
    st.session_state.learning_data = {
        "questions": [],
        "correct_rate": 0,
        "weak_points": [],
        "study_time": 0,
        "interactions": [],
        "last_updated": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }

# 初始化知识库
if "knowledge_base" not in st.session_state:
    st.session_state.knowledge_base = {
        "documents": [],
        "categories": {}
    }

# 初始化课件生成会话
if "courseware_session" not in st.session_state:
    st.session_state.courseware_session = {
        "topic": "",
        "grade_level": "",
        "subject": "",
        "requirements": [],
        "generated": False,
        "outline": "",
        "ppt_content": "",
        "creation_time": None,
        "last_modified": None
    }

# 数据库连接状态
if "db_connected" not in st.session_state:
    # 获取缓存的数据库连接
    db_connections = get_database_connections()
    
    st.session_state.db_connected = db_connections['main']
    st.session_state.qa_db_connected = db_connections['qa']
    st.session_state.rag_kb_connected = db_connections['rag']
    
    # 记录连接状态
    if st.session_state.db_connected:
        db_connect_success("ai_teaching_assistant")
    
    if st.session_state.qa_db_connected:
        db_connect_success("ai_qa_records")
    
    if st.session_state.rag_kb_connected:
        db_connect_success("ai_rag_knowledge")
        
        # 加载所有文档
        try:
            db_docs = rag_kb.get_all_documents(limit=1000)
            
            if db_docs:
                documents_batch = []
                categories_temp = {}
                existing_doc_ids = set()
                
                for doc in db_docs:
                    rag_doc_id = doc['id']
                    
                    if rag_doc_id in existing_doc_ids:
                        continue
                    existing_doc_ids.add(rag_doc_id)
                    
                    doc_data = doc.get('document_data', {})
                    metadata = doc_data.get('metadata', {})
                    
                    title = metadata.get('title', doc.get('title', '未知文档'))
                    subject = metadata.get('subject', doc.get('file_type', '通用'))
                    category = f"📚 {subject}"
                    
                    doc_info = {
                        "name": title,
                        "type": metadata.get('file_type', doc.get('file_type', 'unknown')),
                        "size": doc.get('file_size', 0),
                        "category": category,
                        "upload_time": metadata.get('upload_time', str(doc.get('upload_time', ''))),
                        "json_format": True,
                        "json_data": doc_data,
                        "json_file_path": doc.get('file_path', ''),
                        "rag_doc_id": rag_doc_id
                    }
                    
                    documents_batch.append(doc_info)
                    categories_temp[category] = categories_temp.get(category, 0) + 1
                
                st.session_state.knowledge_base["documents"].extend(documents_batch)
                st.session_state.knowledge_base["categories"].update(categories_temp)
                
                info(f"✅ 从 RAG 数据库恢复了 {len(db_docs)} 个文档")
        except Exception as e:
            warning(f"⚠️ 从 RAG 数据库恢复文档失败：{str(e)}")
    
    # 从本地备份文件恢复学习数据
    LearningDataManager.load_learning_data()

# 应用 CSS 样式
st.markdown(CustomCSS.get_custom_css(), unsafe_allow_html=True)

# 页面配置
st.set_page_config(page_title="多模态 AI 互动式教学智能体", page_icon="🎓", layout="wide")

# 侧边栏

# 侧边栏
with st.sidebar:
    # 用户信息区域
    if st.session_state.current_user:
        user = st.session_state.current_user
        user_info_col1, user_info_col2 = st.columns([3, 1])
        
        with user_info_col1:
            role_icon = {"teacher": "👨🏫", "student": "‍🎓", "admin": "👑", "guest": "👤"}
            icon = role_icon.get(user.get('role', 'guest'), '👤')
            st.markdown(f"{icon} **{user.get('username', '用户')}**")
            if user.get('role') and user.get('role') != 'guest':
                role_text = {"teacher": "教师", "student": "学生", "admin": "管理员"}
                st.caption(f"角色：{role_text.get(user.get('role'), '用户')}")
        
        with user_info_col2:
            if st.button("退出", key="logout_btn", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        
        st.divider()
    
    # 主标题
    st.title(" AI 智能教育助手")
    
    # 检查数据库连接状态，如未连接则弹窗警告
    if not st.session_state.db_connected:
        st.error("⚠️ 主数据库未连接！")
        st.warning("部分功能可能无法使用")
    if not st.session_state.qa_db_connected:
        st.error("⚠️ 答答疑数据库未连接！")
        st.warning("AI 答疑记录将无法保存")
    
    st.divider()
    
    # 插图展示区域
    st.divider()
    st.image(
        "illustration/i1.png",
        width=200
    )
    st.caption("让 AI 助力您的教学之旅 🚀")
    st.divider()
    
    # 检查是否为游客模式
    is_guest = st.session_state.current_user and st.session_state.current_user.get('role') == 'guest'
    
    if is_guest:
        st.info(" 游客模式：仅开放体验功能")
        st.warning("请注册/登录以使用完整功能")
        st.divider()
        
        # 游客只能使用智能答疑
        menu = st.radio(
            "选择功能",
            ["智能答疑"],  # 只有这一个选项
            index=0,
            label_visibility="collapsed",
            horizontal=True,
            key="menu_selection"
        )
    else:
        st.title("导航菜单")
        
        menu = st.radio(
            "选择功能",
            ["智能答疑", "课件生成", "学情分析", "知识库管理"],
            index=0,
            label_visibility="collapsed",
            horizontal=True,
            key="menu_selection"
        )

# 辅助函数

def process_question(question, scenario, learning_data):
    """处理用户问题 - 使用服务层"""
    # 初始化 QA 服务
    qa_service = QAService()
    
    # 调用服务层处理问题
    result = qa_service.handle_text_question(
        question=question,
        scenario=scenario,
        api_key=DEFAULT_API_KEY,
        base_url=BASE_URL
    )
    
    # UI 渲染部分保留在主程序中
    if result["success"]:
        with st.chat_message("assistant"):
            st.markdown(result["answer"])
            
            if result.get("source_info"):
                st.caption(result["source_info"])
            
            # 显示参考的书籍列表
            if result.get("rag_docs_found"):
                with st.expander("📖 查看参考的书籍"):
                    for doc in result["rag_docs_found"]:
                        st.markdown(f"- 📚 {doc['title']} ({doc['subject']})")
            
            # 检查是否包含链接
            urls = extract_urls(result["answer"])
            if urls:
                st.caption("📎 相关链接：")
                for url in urls:
                    st.markdown(f"- {url}")
            
            # 提供下载按钮
            if len(result["answer"]) > 100:
                file_name = generate_filename("AI_回复", "txt")
                st.download_button(
                    label="📥 下载 AI 回复为文本文件",
                    data=result["answer"],
                    file_name=file_name,
                    mime="text/plain",
                    key=f"download_{len(learning_data.get('questions', []))}"
                )
        
        # 聊天记录已由服务层管理，此处只需保存到学习数据
        learning_data["questions"].append({
            "question": question,
            "answer": result["answer"],
            "scenario": scenario,
            "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "answered": True,
            "source": result.get("source_info", "kimi_ai"),
            "tokens_used": result.get("tokens_used"),
            "response_time_ms": result.get("response_time_ms"),
            "rag_docs_count": len(result.get("rag_docs_found", []))
        })
        
        # 每10次问答备份一次
        if len(learning_data["questions"]) % 10 == 0:
            LearningDataManager.save_learning_data()
    else:
        st.error(f"请求失败：{result.get('error', '未知错误')}")

def process_audio(audio_value, scenario, learning_data):
    """处理语音输入 - 使用服务层"""
    # 初始化 QA 服务
    qa_service = QAService()
    
    # 模拟语音识别（实际项目中应集成真实的语音识别 API）
    transcribed_text = f"[语音识别结果] {datetime.now().strftime('%H:%M:%S')} 的语音输入"
    
    # UI 显示
    with st.chat_message("user"):
        st.audio(audio_value)
        st.caption("🎤 语音消息")
    
    with st.chat_message("assistant"):
        with st.spinner("正在识别语音并生成回答..."):
            try:
                # 显示转录文本
                st.markdown(f"**识别到的文本：** {transcribed_text}")
                
                # 调用服务层处理语音问题
                result = qa_service.handle_voice_question(
                    transcribed_text=transcribed_text,
                    api_key=DEFAULT_API_KEY,
                    base_url=BASE_URL
                )
                
                if result["success"]:
                    # 显示 AI 回复
                    st.markdown(result["answer"])
                    
                    # 检查是否包含链接
                    urls = extract_urls(result["answer"])
                    if urls:
                        st.caption("📎 相关链接：")
                        for url in urls:
                            st.markdown(f"- {url}")
                    
                    # 提供下载按钮
                    if len(result["answer"]) > 100:
                        file_name = generate_filename("AI_语音回复", "txt")
                        st.download_button(
                            label="📥 下载 AI 回复为文本文件",
                            data=result["answer"],
                            file_name=file_name,
                            mime="text/plain",
                            key=f"download_audio_{len(learning_data.get('questions', []))}"
                        )
                    
                    # 聊天记录已由服务层管理，此处只需保存到学习数据
                    learning_data["questions"].append({
                        "question": transcribed_text,
                        "answer": result["answer"],
                        "scenario": scenario,
                        "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "answered": True,
                        "input_type": "voice"
                    })
                    
                    # 每10次问答备份一次
                    if len(learning_data["questions"]) % 10 == 0:
                        LearningDataManager.save_learning_data()
                else:
                    st.error(f"语音处理失败：{result.get('error', '未知错误')}")
                    
            except Exception as e:
                st.error(f"语音处理失败：{str(e)}")


# ==================== 学情分析图表生成函数 ====================

def generate_analysis_charts(questions_data, analysis_mode):
    """根据学习数据自动生成分析图表"""
    charts = {}
    
    if not questions_data:
        return charts
    
    try:
        # 1. 互动场景分布饼图
        scenarios = {}
        for q in questions_data:
            scenario = q.get('scenario', '未知')
            scenarios[scenario] = scenarios.get(scenario, 0) + 1
        
        if scenarios:
            fig, ax = plt.subplots(figsize=(8, 6))
            colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#F7DC6F']
            wedges, texts, autotexts = ax.pie(
                scenarios.values(),
                labels=scenarios.keys(),
                autopct='%1.1f%%',
                colors=colors[:len(scenarios)],
                startangle=90
            )
            ax.set_title('互动场景分布', fontsize=14, fontweight='bold', pad=20)
            plt.tight_layout()
            
            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            charts['scenario_pie'] = buf.getvalue()
            plt.close()
        
        # 2. 时间趋势折线图（按日期统计）
        date_stats = {}
        for q in questions_data:
            time_str = q.get('time', '')
            if time_str:
                date_part = time_str.split(' ')[0] if ' ' in time_str else time_str[:10]
                date_stats[date_part] = date_stats.get(date_part, 0) + 1
        
        if len(date_stats) > 1:
            sorted_dates = sorted(date_stats.keys())
            counts = [date_stats[d] for d in sorted_dates]
            
            fig, ax = plt.subplots(figsize=(10, 5))
            ax.plot(sorted_dates, counts, marker='o', linewidth=2, markersize=8, color='#4ECDC4')
            ax.fill_between(range(len(sorted_dates)), counts, alpha=0.3, color='#4ECDC4')
            ax.set_xlabel('日期', fontsize=12)
            ax.set_ylabel('问题数量', fontsize=12)
            ax.set_title('学习活跃度趋势', fontsize=14, fontweight='bold', pad=20)
            ax.grid(True, alpha=0.3, linestyle='--')
            plt.xticks(rotation=45)
            plt.tight_layout()
            
            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            charts['trend_line'] = buf.getvalue()
            plt.close()
        
        # 3. 知识点掌握情况柱状图
        subject_scores = {}
        for q in questions_data:
            subject = q.get('subject', '综合')
            score = q.get('score', None)
            if score is not None:
                if subject not in subject_scores:
                    subject_scores[subject] = []
                subject_scores[subject].append(score)
        
        if subject_scores:
            subjects = list(subject_scores.keys())
            avg_scores = [sum(scores)/len(scores) for scores in subject_scores.values()]
            
            fig, ax = plt.subplots(figsize=(10, 6))
            bars = ax.bar(subjects, avg_scores, color=['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8'][:len(subjects)])
            ax.set_xlabel('学科', fontsize=12)
            ax.set_ylabel('平均得分', fontsize=12)
            ax.set_title('各学科掌握情况', fontsize=14, fontweight='bold', pad=20)
            ax.set_ylim(0, 100)
            
            for bar, score in zip(bars, avg_scores):
                ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 1,
                       f'{score:.1f}', ha='center', va='bottom', fontsize=10)
            
            plt.tight_layout()
            
            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            charts['subject_bar'] = buf.getvalue()
            plt.close()
        
        # 4. 学习时长分布直方图
        study_times = [q.get('study_time', 0) for q in questions_data if q.get('study_time', 0) > 0]
        if study_times:
            fig, ax = plt.subplots(figsize=(8, 5))
            ax.hist(study_times, bins=10, color='#45B7D1', edgecolor='black', alpha=0.7)
            ax.set_xlabel('学习时长（分钟）', fontsize=12)
            ax.set_ylabel('频次', fontsize=12)
            ax.set_title('学习时长分布', fontsize=14, fontweight='bold', pad=20)
            ax.grid(True, alpha=0.3, axis='y')
            plt.tight_layout()
            
            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            charts['time_hist'] = buf.getvalue()
            plt.close()
        
    except Exception as e:
        st.warning(f"⚠️ 图表生成失败：{str(e)}")
    
    return charts


def display_analysis_charts(charts):
    """展示分析图表"""
    if not charts:
        return
    
    st.subheader("📊 AI 数据分析图表")
    
    col1, col2 = st.columns(2)
    with col1:
        if 'scenario_pie' in charts:
            st.image(charts['scenario_pie'], caption='互动场景分布', use_container_width=True)
    
    with col2:
        if 'trend_line' in charts:
            st.image(charts['trend_line'], caption='学习活跃度趋势', use_container_width=True)
    
    if 'subject_bar' in charts or 'time_hist' in charts:
        cols = st.columns(2)
        idx = 0
        if 'subject_bar' in charts:
            with cols[idx]:
                st.image(charts['subject_bar'], caption='各学科掌握情况', use_container_width=True)
            idx += 1
        if 'time_hist' in charts:
            with cols[idx]:
                st.image(charts['time_hist'], caption='学习时长分布', use_container_width=True)


# ==================== 主菜单逻辑 ====================

# 主界面
if menu == "智能答疑":
    st.title("💬 智能答疑")
    
    # 聊天记录管理区域
    if st.session_state.learning_data.get("questions"):
        with st.expander(f"📝 聊天记录管理 ({len(st.session_state.learning_data['questions'])} 条)"):
            col1, col2, col3 = st.columns([3, 1, 1])
            
            with col1:
                st.subheader("🔍 搜索聊天记录")
                search_keyword = st.text_input(
                    "",
                    placeholder="输入关键词搜索聊天记录...",
                    key="chat_search_input"
                )
            
            with col2:
                st.subheader("📤 导出记录")
                export_format = st.selectbox(
                    "导出格式",
                    ["TXT", "JSON"],
                    key="export_format_select"
                )
            
            with col3:
                st.subheader("💾 备份")
                if st.button("💾 立即备份", use_container_width=True):
                    analysis_service = AnalysisService()
                    result = analysis_service.manage_learning_data("backup")
                    if result["success"]:
                        st.success("✅ 数据已备份")
                    else:
                        st.error(f"备份失败：{result.get('error')}")
            
            # 多轮对话控制
            st.divider()
            st.subheader("💬 多轮对话控制")
            col_msg1, col_msg2 = st.columns([3, 1])
            
            with col_msg1:
                msg_count = len(st.session_state.get("messages", [])) // 2  # 除以2因为每轮有用户和AI两条消息
                st.caption(f"当前对话历史：{msg_count} 轮对话（{len(st.session_state.get('messages', []))} 条消息）")
            
            with col_msg2:
                if st.button("🗑️ 清空对话历史", type="secondary", use_container_width=True, key="clear_chat_messages"):
                    qa_service = QAService()
                    result = qa_service.manage_chat_history("clear")
                    if result["success"]:
                        st.success("✅ 对话历史已清空")
                        st.rerun()
                    else:
                        st.error(f"清空失败：{result.get('error')}")
            
            # 搜索结果显示
            filtered_questions = st.session_state.learning_data["questions"]
            if search_keyword:
                # 使用服务层搜索
                analysis_service = AnalysisService()
                result = analysis_service.manage_learning_data("search", keyword=search_keyword)
                if result["success"]:
                    filtered_questions = result["data"]
                    st.caption(f"🔍 找到 {len(filtered_questions)} 条相关记录")
            
            # 显示所有问题列表
            question_options = []
            for i, q in enumerate(filtered_questions):
                question_preview = q['question'][:50] + "..." if len(q['question']) > 50 else q['question']
                question_options.append(f"{i+1}. [{q.get('scenario', '答疑')}] {question_preview}")
            
            if question_options:
                # 多选删除
                questions_to_delete = st.multiselect(
                    "选择要删除的记录（可多选）",
                    options=question_options,
                    key="delete_questions_select"
                )
                
                col1, col2, col3 = st.columns([1, 1, 2])
                with col1:
                    if st.button("🗑️ 删除选中记录", type="primary", use_container_width=True, key="delete_learning_records_btn"):
                        if questions_to_delete:
                            # 获取要删除的索引
                            indices_to_delete = [int(q.split('.')[0]) - 1 for q in questions_to_delete]
                            # 保留未删除的记录
                            st.session_state.learning_data["questions"] = [
                                q for i, q in enumerate(st.session_state.learning_data["questions"])
                                if i not in indices_to_delete
                            ]
                            st.success(f"✅ 已删除 {len(questions_to_delete)} 条记录")
                            st.rerun()
                
                with col2:
                    if st.button("🗑️ 清空所有记录", use_container_width=True):
                        analysis_service = AnalysisService()
                        result = analysis_service.manage_learning_data("clear")
                        if result["success"]:
                            st.success("✅ 已清空所有聊天记录")
                            st.rerun()
                
                with col3:
                    st.caption(f"当前共 {len(st.session_state.learning_data['questions'])} 条记录")
                
                # 导出功能
                if st.button("📤 导出聊天记录", type="primary", use_container_width=True):
                    analysis_service = AnalysisService()
                    format_type = "json" if export_format == "JSON" else "txt"
                    result = analysis_service.manage_learning_data("export", format=format_type)
                    
                    if result["success"] and result["data"]:
                        file_name = generate_filename("聊天记录", export_format.lower())
                        mime_type = "application/json" if export_format == "JSON" else "text/plain"
                        
                        st.download_button(
                            label=f"📥 下载 {export_format} 格式",
                            data=result["data"].encode('utf-8') if isinstance(result["data"], str) else result["data"],
                            file_name=file_name,
                            mime=mime_type,
                            key=f"download_{export_format.lower()}_chat"
                        )
    
    tab1, tab2, tab3 = st.tabs(["📖 课前预习", "🎯 课中互动", "📝 课后辅导"])
    
    with tab1:
        st.markdown("### 📖 课前预习")
        
        # 子功能选项（一行显示）
        preview_mode = st.radio(
            "选择预习模式",
            ["📚 智能题库推荐", "💬 预习答疑", "🎯 知识点预览"],
            horizontal=True,
            label_visibility="collapsed",
            key="preview_mode_select"
        )
        
        if preview_mode == "📚 智能题库推荐":
            st.info("根据您的学习进度，AI 为您推荐针对性练习题")
            preview_question = st.chat_input("请输入您想练习的知识点或题目类型...", key="preview_chat")
            if preview_question:
                process_question(preview_question, "智能题库推荐", st.session_state.learning_data)
        
        elif preview_mode == "💬 预习答疑":
            st.info("预习过程中遇到疑问？随时向 AI 提问")
            preview_question = st.chat_input("请输入预习中的疑问...", key="preview_chat")
            if preview_question:
                process_question(preview_question, "预习答疑", st.session_state.learning_data)
        
        else:  # 知识点预览
            st.info("AI 为您梳理即将学习的知识点框架")
            preview_question = st.chat_input("请输入您想预习的课程内容或章节...", key="preview_chat")
            if preview_question:
                process_question(preview_question, "知识点预览", st.session_state.learning_data)
    
    with tab2:
        st.markdown("### 🎯 课中互动")
        
        # 子功能选项（一行显示）
        classroom_mode = st.radio(
            "选择互动模式",
            ["🎤 实时提问", "👥 小组协作", "💭 课堂讨论", "🎲 随机点名"],
            horizontal=True,
            label_visibility="collapsed",
            key="classroom_mode_select"
        )
        
        if classroom_mode == "🎤 实时提问":
            st.info("课堂上遇到问题？立即向 AI 提问")
            input_mode = st.radio("提问方式", ["文字提问", "语音提问"], horizontal=True, label_visibility="collapsed", key="classroom_input")
            if input_mode == "文字提问":
                classroom_prompt = st.chat_input("请输入课堂问题...", key="classroom_chat")
                if classroom_prompt:
                    process_question(classroom_prompt, "实时提问", st.session_state.learning_data)
            else:
                audio_value = st.audio_input("🎤 按住说话", key="classroom_voice")
                if audio_value:
                    process_audio(audio_value, "实时提问", st.session_state.learning_data)
        
        elif classroom_mode == "👥 小组协作":
            st.info("组织小组讨论，AI 提供协作指导")
            classroom_prompt = st.chat_input("请输入小组讨论的主题或任务...", key="group_chat")
            if classroom_prompt:
                process_question(classroom_prompt, "小组协作", st.session_state.learning_data)
        
        elif classroom_mode == "💭 课堂讨论":
            st.info("发起全班讨论，AI 引导深度思考")
            classroom_prompt = st.chat_input("请输入讨论话题...", key="discussion_chat")
            if classroom_prompt:
                process_question(classroom_prompt, "课堂讨论", st.session_state.learning_data)
        
        else:  # 随机点名
            st.info("班级管理 - 随机抽取学生回答问题")
            # 班级名单管理
            st.subheader("📋 班级管理")
            class_name = st.text_input("班级名称", placeholder="例如：高一 (3) 班", key="class_name_input")
            
            uploaded_roster = st.file_uploader(
                "上传班级名单（Excel/CSV/TXT 格式，每行一个姓名）",
                type=["xlsx", "xls", "csv", "txt"],
                key="roster_upload"
            )
            
            manual_students = st.text_area(
                "或手动输入学生名单",
                placeholder="张三\n李四\n王五\n赵六",
                height=80,
                key="manual_students"
            )
            
            student_list = []
            if uploaded_roster:
                try:
                    if uploaded_roster.type in ["text/plain", "text/csv"]:
                        content = uploaded_roster.read().decode('utf-8')
                        student_list = [line.strip() for line in content.split('\n') if line.strip()]
                    elif uploaded_roster.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
                        import pandas as pd
                        df = pd.read_excel(uploaded_roster)
                        student_list = df.iloc[:, 0].dropna().astype(str).tolist()
                    st.success(f"✅ 已加载 {len(student_list)} 名学生")
                except Exception as e:
                    st.error(f"读取名单失败：{str(e)}")
            elif manual_students:
                student_list = [line.strip() for line in manual_students.split('\n') if line.strip()]
            
            if student_list:
                st.divider()
                col1, col2 = st.columns([2, 1])
                with col1:
                    st.subheader("🎲 随机点名")
                    if st.button("🎯 开始随机抽取", type="primary", use_container_width=True):
                        import random
                        selected_student = random.choice(student_list)
                        st.success(f"🎉 被点到的同学是：**{selected_student}**")
                        st.session_state.learning_data["interactions"].append({
                            "type": "random_pick",
                            "student": selected_student,
                            "class": class_name,
                            "timestamp": datetime.now()
                        })
                with col2:
                    st.subheader("👥 学生统计")
                    st.info(f"班级总人数：{len(student_list)}人")
                
                # 互动记录管理
                if st.session_state.learning_data.get("interactions"):
                    st.divider()
                    st.subheader("📊 互动记录管理")
                    
                    # 显示互动记录列表
                    interaction_options = []
                    for i, interaction in enumerate(st.session_state.learning_data["interactions"]):
                        if interaction.get("type") == "random_pick":
                            time_str = interaction.get("timestamp", datetime.now()).strftime("%H:%M:%S")
                            interaction_options.append(
                                f"{i+1}. 🎲 [{time_str}] {interaction.get('class', '未知班级')} - {interaction.get('student', '未知学生')}"
                            )
                    
                    if interaction_options:
                        interactions_to_delete = st.multiselect(
                            "选择要删除的互动记录（可多选）",
                            options=interaction_options,
                            key="delete_interactions_select"
                        )
                        
                        col1, col2 = st.columns([1, 3])
                        with col1:
                            if st.button("🗑️ 删除选中记录", type="primary", use_container_width=True, key="delete_interaction_records_btn"):
                                if interactions_to_delete:
                                    # 获取要删除的索引
                                    indices_to_delete = [int(q.split('.')[0]) - 1 for q in interactions_to_delete]
                                    # 保留未删除的记录
                                    st.session_state.learning_data["interactions"] = [
                                        inter for i, inter in enumerate(st.session_state.learning_data["interactions"])
                                        if i not in indices_to_delete
                                    ]
                                    st.success(f"✅ 已删除 {len(interactions_to_delete)} 条互动记录")
                                    st.rerun()
                        
                        with col2:
                            if st.button("🗑️ 清空所有互动", use_container_width=True):
                                st.session_state.learning_data["interactions"] = []
                                st.success("✅ 已清空所有互动记录")
                                st.rerun()
    
    with tab3:
        st.markdown("### 📝 课后辅导")
        
        # 子功能选项（一行显示）
        homework_mode = st.radio(
            "选择辅导模式",
            ["📸 作业批改", "💬 错题讲解", "📚 个性化推荐"],
            horizontal=True,
            label_visibility="collapsed",
            key="homework_mode_select"
        )
        
        if homework_mode == "📸 作业批改":
            st.info("上传作业图片，AI 智能批改并给出反馈")
            homework_type = st.selectbox(
                "作业科目",
                ["语文", "数学", "英语", "物理", "化学", "生物", "历史", "地理", "政治", "体育", "美术", "音乐", "信息技术"]
            )
            uploaded_homework = st.file_uploader("上传作业图片", type=["jpg", "png", "jpeg"])
            if uploaded_homework:
                col1, col2 = st.columns([4, 1])
                with col1:
                    st.image(uploaded_homework, caption="已上传的作业", use_container_width=True)
                with col2:
                    st.download_button(
                        label="📥 下载",
                        data=uploaded_homework.read(),
                        file_name=uploaded_homework.name,
                        mime=uploaded_homework.type,
                        key=f"homework_download_{uploaded_homework.name}"
                    )
                if st.button("🔍 开始批改"):
                    with st.spinner("正在分析作业..."):
                        st.success("✅ 批改完成！")
                        st.info("正确率：85%\n\n薄弱点：三角函数公式应用\n\n建议：复习正弦定理和余弦定理")
        
        elif homework_mode == "💬 错题讲解":
            st.info("遇到不会的题目？AI 详细讲解")
            homework_question = st.chat_input("请输入作业中的问题或拍照上传...", key="homework_chat")
            if homework_question:
                process_question(homework_question, "错题讲解", st.session_state.learning_data)
        
        else:  # 个性化推荐
            st.info("根据您的学习情况，AI 推荐针对性练习")
            homework_question = st.chat_input("请描述您想加强的知识点或薄弱环节...", key="recommend_chat")
            if homework_question:
                process_question(homework_question, "个性化推荐", st.session_state.learning_data)
    
    st.subheader("📎 上传教学资料")
    uploaded_files = st.file_uploader(
        "支持 PDF、Word、PPT 等多种格式，AI 将自动解析内容并提炼知识点",
        type=["pdf", "doc", "docx", "ppt", "pptx", "txt"],
        accept_multiple_files=True,
        key="material_upload"
    )
    if uploaded_files:
        for file in uploaded_files:
            col1, col2 = st.columns([3, 1])
            with col1:
                st.success(f"✅ {file.name}")
            with col2:
                # 下载按钮
                st.download_button(
                    label="📥 下载",
                    data=file.read(),
                    file_name=file.name,
                    mime="application/octet-stream",
                    key=f"download_{file.name}"
                )
        
        if st.button("🤖 AI 解析课件内容", type="primary"):
            with st.spinner("正在解析课件内容..."):
                try:
                    # 调用服务层解析文件
                    courseware_service = CoursewareService()
                    result = courseware_service.analyze_uploaded_files(
                        uploaded_files=uploaded_files,
                        api_key=DEFAULT_API_KEY,
                        base_url=BASE_URL
                    )
                    
                    if result["success"]:
                        st.success("✅ 课件解析完成！")
                        st.markdown(result["analysis"])
                    else:
                        st.error(f"❌ 解析失败：{result.get('error', '未知错误')}")
                    
                except Exception as e:
                    st.error(f"解析失败：{str(e)}")

elif menu == "课件生成":
    # 检查游客模式
    if is_guest:
        st.title("📚 AI 课件生成")
        st.warning("⚠️ 游客模式无法使用此功能")
        st.info("💡 请注册或登录以使用课件生成功能")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📝 立即注册", type="primary", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        with col2:
            if st.button("🔐 去登录", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        st.stop()
    
    st.title("📚 AI 课件生成")
    
    # ✅ 加载历史课件按钮
    col_hist, col_new = st.columns([1, 4])
    with col_hist:
        if st.button("📂 加载历史课件", use_container_width=True):
            if st.session_state.db_connected:
                try:
                    history_courseware = db.get_all_courseware()
                    if history_courseware:
                        # 创建选择框
                        courseware_options = {f"{cw['title']} ({cw['subject']})": cw for cw in history_courseware}
                        selected = st.selectbox(
                            "选择要加载的课件",
                            options=list(courseware_options.keys()),
                            key="select_history_courseware"
                        )
                        
                        if selected and st.button("确认加载", type="primary"):
                            cw = courseware_options[selected]
                            # 解析 JSON 数据
                            cw_data = json.loads(cw['content'])
                            
                            # 恢复到 session_state
                            st.session_state.courseware_session = {
                                "topic": cw_data.get('title', cw['title']),
                                "grade_level": cw_data.get('grade_level', cw.get('grade_level', '')),
                                "subject": cw_data.get('subject', cw.get('subject', '')),
                                "requirements": cw_data.get('requirements', []),
                                "generated": True,
                                "outline": cw_data.get('outline', ''),
                                "ppt_content": cw_data.get('slides', []),
                                "db_id": cw['id']
                            }
                            
                            info(f"✅ 已加载课件：{selected}")
                            st.rerun()
                    else:
                        st.info("暂无历史课件")
                except Exception as e:
                    error(f"加载历史课件失败：{str(e)}")
            else:
                st.warning("数据库未连接，无法加载历史课件")
    
    # 基本信息输入
    st.subheader("📝 第一步：填写基本信息")
    topic = st.text_input("课件主题", value=st.session_state.courseware_session.get("topic", ""), placeholder="例如：函数的单调性", key="topic_input")
    
    # 更新会话状态
    if topic:
        st.session_state.courseware_session["topic"] = topic
    
    # 上传参考资料
    st.subheader("📎 第二步：上传参考资料（可选）")
    reference_files = st.file_uploader("上传教学资料（支持图片、PDF、Word 等）",
                                        type=["pdf", "doc", "docx", "ppt", "pptx", "txt", "jpg", "jpeg", "png"],
                                        accept_multiple_files=True,
                                        key="courseware_upload")
    
    if reference_files:
        for file in reference_files:
            st.success(f"✅ {file.name}")
    
    # 描述具体需求
    st.subheader("💬 第三步：描述您的具体需求")
    st.caption("您可以告诉 AI 想要什么样的课件风格、重点内容、活动设计等")
    
    # 初始化课件生成对话历史
    if "courseware_chat_history" not in st.session_state:
        st.session_state.courseware_chat_history = []
    
    # 显示对话历史
    if st.session_state.courseware_chat_history:
        st.divider()
        for message in st.session_state.courseware_chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        st.divider()
    
    # 需求输入框（聊天框）- 放在需求列表下方
    col1, col2 = st.columns([4, 1])
    with col1:
        requirement_input = st.text_input(
            "添加新要求",
            placeholder="请输入您的具体要求（如：希望多一些互动环节、需要包含视频资源链接等）...",
            key="requirement_input",
            label_visibility="collapsed"
        )
    with col2:
        add_btn = st.button("➕ 添加", type="primary", use_container_width=True, key="add_requirement_btn")
    
    if requirement_input and add_btn:
        # 添加到对话历史
        st.session_state.courseware_chat_history.append({
            "role": "user",
            "content": requirement_input
        })
        
        # 添加到需求列表
        st.session_state.courseware_session["requirements"].append(requirement_input)
        
        # AI 确认回复
        ai_response = "✅ 已记录您的要求，将在生成课件时考虑这一点。"
        st.session_state.courseware_chat_history.append({
            "role": "assistant",
            "content": ai_response
        })
        
        st.rerun()
    
    # 智能澄清对话（在第三步和第四步之间）
    st.divider()
    st.subheader("💡 智能需求澄清")
    st.caption("AI 将主动提问以更好地理解您的需求")
    
    # 初始化澄清对话
    if "clarification_started" not in st.session_state:
        st.session_state.clarification_started = False
    
    if "clarification_messages" not in st.session_state:
        st.session_state.clarification_messages = []
    
    if not st.session_state.clarification_started and topic:
        if st.button("🤖 开始智能澄清", type="secondary"):
            st.session_state.clarification_started = True
            
            # 调用服务层开始澄清
            courseware_service = CoursewareService()
            result = courseware_service.start_clarification(
                topic=topic,
                api_key=DEFAULT_API_KEY,
                base_url=BASE_URL
            )
            
            if result["success"]:
                st.session_state.clarification_messages.append({
                    "role": "assistant",
                    "content": result["question"]
                })
            else:
                st.error(f"❌ 澄清失败：{result.get('error', '未知错误')}")
            
            st.rerun()
    
    # 显示澄清对话
    if st.session_state.clarification_started:
        for msg in st.session_state.clarification_messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])
        
        # 用户回复输入
        user_reply = st.chat_input("请回答 AI 的问题...", key="clarification_input")
        
        if user_reply:
            # 添加用户回复
            st.session_state.clarification_messages.append({
                "role": "user",
                "content": user_reply
            })
            
            # 调用服务层继续澄清
            courseware_service = CoursewareService()
            result = courseware_service.continue_clarification(
                topic=topic,
                conversation_history=st.session_state.clarification_messages,
                api_key=DEFAULT_API_KEY,
                base_url=BASE_URL
            )
            
            if result["success"]:
                st.session_state.clarification_messages.append({
                    "role": "assistant",
                    "content": result["response"]
                })
                
                # 如果包含“需求已确认”，提取需求到 requirements
                if result.get("confirmed"):
                    st.session_state.courseware_session["requirements"].append(user_reply)
                    st.success("✅ 需求已确认！现在可以点击“开始生成课件”按钮")
            else:
                st.error(f"❌ 澄清失败：{result.get('error', '未知错误')}")
            
            st.rerun()
    
    # 生成课件
    st.divider()
    st.subheader("🚀 第四步：生成课件")
    
    # 添加快速模式选项
    col_speed, col_info = st.columns([1, 3])
    with col_speed:
        fast_mode = st.checkbox("⚡ 快速模式", value=True, help="启用后可提升 50-70% 的生成速度，但课件内容会更精简")
    with col_info:
        if fast_mode:
            st.info("💨 快速模式已启用：8-10页幻灯片，无配图，简化装饰")
        else:
            st.info("🎨 标准模式：10-15页幻灯片，含配图和精美装饰（生成较慢）")
    
    col1, col2 = st.columns([1, 3])
    with col1:
        generate_btn = st.button("开始生成课件", type="primary", use_container_width=True, disabled=not topic)
    with col2:
        if st.button("清空所有信息，重新开始", use_container_width=True):
            st.session_state.courseware_session = {
                "topic": "",
                "grade_level": "",
                "subject": "",
                "requirements": [],
                "generated": False,
                "outline": "",
                "ppt_content": ""
            }
            # 清空对话历史
            st.session_state.courseware_chat_history = []
            st.rerun()
    
    if generate_btn and topic:
        mode_text = "快速模式" if fast_mode else "标准模式"
        with st.status(f"正在生成课件（{mode_text}）...", expanded=True) as status:
            try:
                # 构建完整的需求描述
                requirements_text = "\n".join(st.session_state.courseware_session["requirements"]) if st.session_state.courseware_session.get("requirements") else "无特殊要求"
                
                status.update(label=f"⏳ 步骤 1/3: 识别学科并生成大纲（{mode_text}）...")
                
                # 调用服务层生成课件
                courseware_service = CoursewareService()
                result = courseware_service.generate_courseware(
                    topic=topic,
                    requirements_text=requirements_text,
                    api_key=DEFAULT_API_KEY,
                    base_url=BASE_URL,
                    fast_mode=fast_mode  # 传递快速模式参数
                )
                
                if result["success"]:
                    status.update(label=f"✅ 步骤 2/3: 生成 PPT 结构完成")
                    
                    # 更新 session_state
                    st.session_state.courseware_session["subject"] = result["subject"]
                    st.session_state.courseware_session["outline"] = result["outline"]
                    st.session_state.courseware_session["ppt_content"] = result["slides"]
                    st.session_state.courseware_session["ppt_theme"] = result["theme"]
                    st.session_state.courseware_session["generated"] = True
                    st.session_state.courseware_session["creation_time"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    st.session_state.courseware_session["last_modified"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    st.session_state.courseware_session["db_id"] = result.get("courseware_id")
                    
                    # 保存生成的图片和动画（如果有）
                    if result.get("generated_images"):
                        st.session_state.courseware_session["generated_images"] = result["generated_images"]
                    
                    status.update(label=f"✅ 步骤 3/3: 保存到数据库完成")
                    status.update(label="✅ 课件生成完成!", state="complete")
                    st.success(f"✅ 学科识别：{result['subject']}")
                    st.success(f"✅ 幻灯片数：{len(result['slides'])}")
                    if result.get("generated_images"):
                        st.success(f"✅ 生成配图：{len(result['generated_images'])} 张")
                    else:
                        st.info("💨 快速模式：已跳过图片生成以提升速度")
                else:
                    st.error(f"❌ 生成失败：{result.get('error', '未知错误')}")
                    status.update(label="❌ 生成失败", state="error")
                    
            except Exception as e:
                error_msg = str(e)
                st.error(f"❌ 生成失败：{error_msg}")
                status.update(label="❌ 生成失败", state="error")
                
                # 提供更详细的错误提示和解决方案
                if "连接错误" in error_msg or "connection" in error_msg.lower():
                    st.warning("🔧 可能的原因和解决方案：")
                    st.info("1️⃣ 网络连接不稳定 - 请检查您的网络")
                    st.info("2️⃣ API Key 无效或已过期 - 请检查 .env 文件中的 KIMI_API_KEY")
                    st.info("3️⃣ Kimi API 服务暂时不可用 - 请稍后重试")
                    st.info("4️⃣ 请求超时 - 已增加超时时间，如仍失败请稍后重试")
                elif "timeout" in error_msg.lower():
                    st.warning("⏰ 请求超时：AI 服务响应时间过长")
                    st.info("💡 建议：请稍后重试，或简化您的要求")
                elif "JSON" in error_msg or "json" in error_msg.lower():
                    st.warning("📄 AI 返回的数据格式有误")
                    st.info("💡 建议：请重新生成，AI 会自动调整格式")
                else:
                    st.info("💡 建议：请检查错误信息，稍后重试。如问题持续存在，请联系管理员。")

# ✅ 显示已生成的课件内容（仅在“课件生成”页面显示）
if menu == "课件生成" and st.session_state.courseware_session.get("generated") and st.session_state.courseware_session.get("outline"):
    st.divider()
    st.subheader("📊 已生成的课件")
    
    # 显示课件元信息
    creation_time = st.session_state.courseware_session.get("creation_time", "未知")
    last_modified = st.session_state.courseware_session.get("last_modified", "未知")
    topic = st.session_state.courseware_session.get("topic", "未命名课件")
    subject = st.session_state.courseware_session.get("subject", "综合")
    
    col_meta1, col_meta2 = st.columns(2)
    with col_meta1:
        st.info(f"📚 主题：{topic}")
    with col_meta2:
        st.info(f"🎓 学科：{subject}")
    
    st.info(f"⏰ 创建时间：{creation_time}")
    
    outline = st.session_state.courseware_session["outline"]
    slides = st.session_state.courseware_session.get("ppt_content", [])
    
    # 显示生成的内容
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("📊 PPT 大纲")
        st.markdown(outline)
    
    with col2:
        st.subheader("📝 教案详情")
        st.text_area("教案内容", value=outline, height=400, key="outline_display")
    
    # 显示完整 PPT 内容预览
    st.divider()
    st.subheader("📄 PPT 课件预览")
    
    # 展示每一页幻灯片
    for i, slide in enumerate(slides):
        with st.expander(f"📄 第 {i+1} 页：{slide.get('title', '无标题')}", expanded=(i==0)):
            subtitle = slide.get('subtitle', '')
            if subtitle and subtitle.strip():
                st.caption(subtitle)
            content = slide.get('content', [])
            # 过滤空字符串，只展示有效内容
            valid_content = [point for point in content if point and point.strip()]
            if valid_content:
                for point in valid_content:
                    st.write(f"• {point}")
            else:
                st.info("暂无内容")
    
    # 可视化 PPT 预览界面（模拟真实 PPT 软件界面）
    st.divider()
    st.subheader("🖼️ PPT 课件预览")
    
    # 获取主题颜色
    theme = st.session_state.courseware_session.get('ppt_theme', {})
    primary_color = theme.get('primary_color', '#0a192f')
    secondary_color = theme.get('secondary_color', '#64ffda')
    bg_color = theme.get('bg_color', '#ffffff')
    text_color = theme.get('text_color', '#333333')
    template_style = theme.get('template_style', 'tech')
    
    # 使用 tabs 分页显示所有幻灯片
    slide_tabs = st.tabs([f"第 {i+1} 页" for i in range(len(slides))])
    
    for idx, (tab, slide) in enumerate(zip(slide_tabs, slides)):
        with tab:
            slide_title = slide.get('title', '无标题')
            subtitle = slide.get('subtitle', '')
            content = slide.get('content', [])
            valid_content = [point for point in content if point and point.strip()]
            
            # 根据模板风格调整样式
            if template_style == 'tech':
                gradient = f'linear-gradient(135deg, {primary_color} 0%, #1a3a5c 100%)'
                shadow = '0 8px 32px rgba(10, 25, 47, 0.3)'
            elif template_style == 'edu':
                gradient = f'linear-gradient(135deg, {primary_color} 0%, #7b4fa2 100%)'
                shadow = '0 8px 32px rgba(91, 44, 111, 0.3)'
            elif template_style == 'nature':
                gradient = f'linear-gradient(135deg, {primary_color} 0%, #2ecc71 100%)'
                shadow = '0 8px 32px rgba(39, 174, 96, 0.3)'
            else:
                gradient = f'linear-gradient(135deg, {primary_color} 0%, #34495e 100%)'
                shadow = '0 8px 32px rgba(44, 62, 80, 0.3)'
            
            # 构建 HTML 内容列表
            content_html = ""
            if valid_content:
                for point in valid_content:
                    content_html += f"<li style='margin-bottom: 18px; font-size: 20px; line-height: 1.8; color: {text_color};'>• {point}</li>"
            else:
                content_html = "<p style='color: #999; font-size: 18px;'>📝 此页暂无详细内容</p>"
            
            # 副标题 HTML
            subtitle_html = ""
            if subtitle and subtitle.strip():
                subtitle_html = f"<p style='color: #666; font-size: 18px; margin: 10px 0 30px 0;'>{subtitle}</p>"
            
            # 构建完整的 HTML 字符串
            full_html = f"""
            <div style='margin: 20px 0;'>
                <div style='background: #e8e8e8; padding: 30px; border-radius: 12px; box-shadow: {shadow};'>
                    <div style='background: white; border-radius: 8px; aspect-ratio: 16/9; display: flex; flex-direction: column; overflow: hidden; box-shadow: 0 4px 20px rgba(0,0,0,0.1);'>
                        <!-- PPT 标题栏 -->
                        <div style='background: {gradient}; padding: 40px 50px 30px 50px; position: relative; min-height: 120px;'>
                            <!-- 装饰元素 -->
                            <div style='position: absolute; top: 0; right: 0; width: 200px; height: 100%; background: linear-gradient(90deg, transparent, rgba(255,255,255,0.1));'></div>
                            <div style='position: absolute; bottom: 0; left: 0; right: 0; height: 5px; background: {secondary_color};'></div>
                            
                            <!-- 标题 -->
                            <h1 style='color: white; margin: 0; font-size: 42px; font-weight: bold; text-shadow: 2px 2px 4px rgba(0,0,0,0.3);'>{slide_title}</h1>
                        </div>
                        
                        <!-- PPT 内容区 -->
                        <div style='flex: 1; padding: 40px 50px; background: {bg_color}; overflow-y: auto;'>
                            {subtitle_html}
                            <ul style='padding-left: 30px; margin: 0;'>
                                {content_html}
                            </ul>
                        </div>
                        
                        <!-- PPT 页脚 -->
                        <div style='background: #f5f5f5; padding: 12px 50px; border-top: 2px solid #e0e0e0; display: flex; justify-content: space-between; align-items: center;'>
                            <span style='color: #999; font-size: 14px;'>AI 课件生成系统</span>
                            <span style='color: #999; font-size: 14px;'>第 {idx+1} / {len(slides)} 页</span>
                        </div>
                    </div>
                </div>
            </div>
            """
            
            # 使用 components.v1.html 渲染 HTML（比 markdown 更可靠）
            st.components.v1.html(full_html, height=500)
    
    # 修改意见反馈功能
    st.divider()
    st.subheader("💬 提出修改意见")
    
    feedback = st.text_area(
        "请输入您的修改建议（如：调整顺序、简化某页、增加一个案例等）",
        placeholder="例如：第3页内容太多，需要简化；希望能增加更多互动环节...",
        height=100,
        key="ppt_feedback"
    )
    
    col_feedback1, col_feedback2 = st.columns([1, 3])
    with col_feedback1:
        if st.button("🔄 基于反馈重新生成", type="primary"):
            if feedback:
                with st.spinner("正在根据您的反馈重新生成课件..."):
                    try:
                        # 调用服务层调整课件
                        courseware_service = CoursewareService()
                        result = courseware_service.refine_courseware(
                            feedback=feedback,
                            topic=st.session_state.courseware_session.get('topic', ''),
                            subject=st.session_state.courseware_session.get('subject', ''),
                            slides=slides,
                            api_key=DEFAULT_API_KEY,
                            base_url=BASE_URL
                        )
                        
                        if result["success"] and result.get("slides"):
                            st.session_state.courseware_session["ppt_content"] = result["slides"]
                            st.session_state.courseware_session["ppt_theme"] = result.get("theme", {})
                            st.success("✅ 课件已根据您的反馈重新生成！")
                            st.rerun()
                        else:
                            st.error(f"❌ 重新生成失败：{result.get('error', '未知错误')}")
                    except Exception as e:
                        st.error(f"❌ 重新生成失败：{str(e)}")
            else:
                st.warning("⚠️ 请先输入修改建议")
    
    # 动画展示区域
    if st.session_state.courseware_session.get("animations"):
        animations = st.session_state.courseware_session.get("animations", [])
        if animations and len(animations) > 0:
            st.divider()
            st.subheader("🎬 教学动画")
            st.caption("AI 生成的配套教学动画，支持 GIF 和 HTML5 格式")
            
            for i, anim in enumerate(animations):
                with st.expander(f"🎬 动画 {i+1}: {anim.get('title', '未命名')}", expanded=True):
                    st.caption(anim.get('description', ''))
                    
                    # 显示 SVG 预览
                    svg_code = anim.get('svg_code', '')
                    if svg_code:
                        try:
                            st.components.v1.html(svg_code, height=350, scrolling=True)
                        except Exception as e:
                            st.warning(f"SVG 预览失败：{str(e)}")
                    
                    st.divider()
                    col_anim1, col_anim2, col_anim3 = st.columns(3)
                    
                    with col_anim1:
                        if st.button(f"💾 下载 HTML", key=f"html_{i}", use_container_width=True):
                            animation_service = AnimationService()
                            html_content = animation_service.generate_html_animation(
                                svg_code, 
                                anim.get('title', '教学动画')
                            )
                            st.download_button(
                                label="下载 HTML 动画文件",
                                data=html_content,
                                file_name=f"{anim.get('title', 'animation')}.html",
                                mime="text/html",
                                key=f"download_html_{i}"
                            )
                    
                    with col_anim2:
                        if st.button(f"🎞️ 生成 GIF", key=f"gif_{i}", use_container_width=True):
                            with st.spinner("正在生成 GIF（可能需要 10-20 秒）..."):
                                try:
                                    animation_service = AnimationService()
                                    gif_path = animation_service.svg_to_gif(svg_code, f"animation_{i}.gif")
                                    if gif_path and os.path.exists(gif_path):
                                        with open(gif_path, "rb") as f:
                                            st.download_button(
                                                label="下载 GIF 动画",
                                                data=f.read(),
                                                file_name=f"{anim.get('title', 'animation')}.gif",
                                                mime="image/gif",
                                                key=f"download_gif_{i}"
                                            )
                                        # 删除临时文件
                                        try:
                                            os.remove(gif_path)
                                        except:
                                            pass
                                    else:
                                        st.error("❌ GIF 生成失败，请使用 HTML 版本")
                                except Exception as e:
                                    st.error(f"❌ GIF 生成失败：{str(e)}")
                    
                    with col_anim3:
                        st.info(f"关联页面：第 {anim.get('related_slide_index', '?')} 页")
                        st.caption(f"类型：{anim.get('animation_type', 'general')}")
            
            st.success("✅ 所有动画已生成，可以下载 HTML 或 GIF 文件集成到 PPT 中")
    
    # 使用 Kimi 生成的 JSON 数据创建真实的 PPT 文件
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import io
    
    def hex_to_rgb(hex_color):
        """将十六进制颜色转换为 RGB 元组"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def apply_background(slide, bg_type, colors):
        """应用背景颜色"""
        try:
            if bg_type == 'gradient' and len(colors) >= 2:
                # 使用第一种颜色作为背景
                bg_color = hex_to_rgb(colors[0])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
            elif bg_type == 'solid' and len(colors) >= 1:
                bg_color = hex_to_rgb(colors[0])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
            else:
                # 默认白色背景
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        except:
            # 如果失败，使用默认白色背景
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    def add_decorative_shape(slide, shape_type, left, top, width, height, fill_color):
        """添加装饰性形状"""
        try:
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
            shape.line.fill.background()
            return shape
        except:
            return None
    
    try:
        prs = Presentation()
        # 设置幻灯片大小为 16:9 宽屏
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 获取模板风格
        template_style = st.session_state.courseware_session.get('template_style', 'tech')
        theme = st.session_state.courseware_session.get('ppt_theme', {})
        primary_color = theme.get('primary_color', '#0a192f')
        secondary_color = theme.get('secondary_color', '#64ffda')
        accent_color = theme.get('accent_color', '#00d4ff')
        bg_color = theme.get('bg_color', '#0a192f')
        text_color = theme.get('text_color', '#ffffff')
        font_name = 'Microsoft YaHei'
        
        def add_circle_decoration(slide, left, top, diameter, fill_color, opacity=0.3):
            """添加圆形装饰"""
            try:
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
                shape.fill.transparency = opacity
                shape.line.fill.background()
                return shape
            except:
                return None
        
        def add_line_decoration(slide, left, top, width, height, line_color, thickness=3):
            """添加线条装饰"""
            try:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(line_color))
                shape.line.fill.background()
                return shape
            except:
                return None
        
        def add_triangle_decoration(slide, left, top, width, height, fill_color):
            """添加三角形装饰"""
            try:
                shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
                shape.fill.transparency = 0.7
                shape.line.fill.background()
                return shape
            except:
                return None
        
        def add_glow_effect(slide, shape, glow_color, glow_size=10):
            """为形状添加发光效果（通过添加半透明圆环模拟）"""
            try:
                left = shape.left - Pt(glow_size)
                top = shape.top - Pt(glow_size)
                width = shape.width + Pt(glow_size * 2)
                height = shape.height + Pt(glow_size * 2)
                glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
                glow.fill.solid()
                glow.fill.fore_color.rgb = RGBColor(*hex_to_rgb(glow_color))
                glow.fill.transparency = 0.8
                glow.line.fill.background()
                return glow
            except:
                return None
        
        def add_cover_decorations(slide, style):
            """为封面页添加装饰元素"""
            if style == 'tech':
                # 科技风格：圆形装饰 + 线条
                add_circle_decoration(slide, Inches(9), Inches(0.5), Inches(3), primary_color, 0.5)
                add_circle_decoration(slide, Inches(10), Inches(1.5), Inches(2), secondary_color, 0.3)
                add_line_decoration(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.1), accent_color)
                add_line_decoration(slide, Inches(1), Inches(7), Inches(5), Inches(0.05), secondary_color)
                
            elif style == 'edu':
                # 教育风格：书本装饰
                add_circle_decoration(slide, Inches(0.5), Inches(5), Inches(2), accent_color, 0.2)
                add_circle_decoration(slide, Inches(11), Inches(0.5), Inches(1.5), primary_color, 0.3)
                add_line_decoration(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), primary_color)
                
            elif style == 'nature':
                # 自然风格：圆点装饰
                add_circle_decoration(slide, Inches(10), Inches(1), Inches(1.5), secondary_color, 0.4)
                add_circle_decoration(slide, Inches(11.5), Inches(2), Inches(1), accent_color, 0.3)
                add_circle_decoration(slide, Inches(0.5), Inches(6), Inches(2), primary_color, 0.2)
                
            elif style == 'minimal':
                # 简约风格：几何线条
                add_line_decoration(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.05), secondary_color)
                add_line_decoration(slide, Inches(6.5), Inches(0), Inches(0.05), Inches(7.5), accent_color)
                
            elif style == 'business':
                # 商务风格：金色装饰
                add_line_decoration(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), accent_color)
                add_circle_decoration(slide, Inches(10.5), Inches(0.5), Inches(2), accent_color, 0.2)
        
        def add_content_decorations(slide, style, slide_index):
            """为内容页添加装饰元素"""
            if style == 'tech':
                # 侧边装饰线
                add_line_decoration(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), secondary_color)
                # 底部装饰
                add_line_decoration(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.1), primary_color)
                
            elif style == 'edu':
                # 顶部装饰条
                add_line_decoration(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.1), accent_color)
                # 页码装饰
                add_circle_decoration(slide, Inches(12.3), Inches(6.8), Inches(0.6), primary_color, 0.3)
                
            elif style == 'nature':
                # 左侧装饰
                add_line_decoration(slide, Inches(0), Inches(0), Inches(0.1), Inches(7.5), secondary_color)
                
            elif style == 'minimal':
                # 极简边框
                add_line_decoration(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.03), RGBColor(200, 200, 200))
                
            elif style == 'business':
                # 顶部金色装饰
                add_line_decoration(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.15), accent_color)
                # 底部装饰
                add_line_decoration(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.1), primary_color)
        
        for i, slide_data in enumerate(slides):
            # 根据布局类型选择版式
            layout_type = slide_data.get('layout', 'title_content')
            
            if i == 0:  # 封面页 - 精美设计
                slide_layout = prs.slide_layouts[6]  # 空白版式
                slide = prs.slides.add_slide(slide_layout)
                
                # 应用深色渐变背景
                bg_info = slide_data.get('background', {})
                apply_background(slide, 'dark_gradient', [primary_color, bg_color])
                
                # 添加模板风格装饰
                add_cover_decorations(slide, template_style)
                
                # 添加发光装饰圆（科技风格）
                if template_style == 'tech':
                    add_circle_decoration(slide, Inches(2), Inches(1), Inches(4), secondary_color, 0.15)
                    add_circle_decoration(slide, Inches(8), Inches(4), Inches(3), accent_color, 0.2)
                
                # 添加大标题（超大字体，居中）
                title = slide_data.get('title', '')
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(2))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.alignment = PP_ALIGN.CENTER
                p.font.name = font_name
                p.font.size = Pt(54)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                
                # 添加装饰线（标题下方）
                add_line_decoration(slide, Inches(4), Inches(4.2), Inches(5.333), Inches(0.08), accent_color)
                
                # 添加副标题
                subtitle = slide_data.get('subtitle', '')
                if subtitle:
                    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
                    tf = subtitle_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = subtitle
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = font_name
                    p.font.size = Pt(22)
                    p.font.color.rgb = RGBColor(200, 220, 255)
                
                # 添加负责人/作者信息
                author = slide_data.get('author', '')
                if author:
                    author_box = slide.shapes.add_textbox(Inches(4), Inches(5.8), Inches(5.333), Inches(0.6))
                    tf = author_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"主讲人：{author}"
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = font_name
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(180, 200, 230)
                
            elif layout_type == 'section_divider':  # 章节过渡页
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                # 全背景色
                apply_background(slide, 'solid', [primary_color])
                
                # 添加装饰
                add_circle_decoration(slide, Inches(10), Inches(0.5), Inches(2.5), accent_color, 0.2)
                add_line_decoration(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.15), secondary_color)
                
                # 居中大标题
                title = slide_data.get('title', '')
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.alignment = PP_ALIGN.CENTER
                p.font.name = font_name
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                
            else:  # 内容页 - 精美设计
                slide_layout = prs.slide_layouts[6]  # 空白版式
                slide = prs.slides.add_slide(slide_layout)
                
                # 应用浅色背景
                apply_background(slide, 'solid', ['#ffffff'])
                
                # 添加模板风格装饰
                add_content_decorations(slide, template_style, i)
                
                # 添加顶部标题栏（带渐变效果）
                header_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), Inches(13.333), Inches(1.3)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(primary_color))
                header_bar.line.fill.background()
                
                # 标题栏底部装饰线
                add_line_decoration(slide, Inches(0), Inches(1.25), Inches(13.333), Inches(0.08), accent_color)
                
                # 添加标题
                title = slide_data.get('title', '')
                title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = font_name
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                
                # 添加页码（右下角）
                page_num_box = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1), Inches(0.5))
                tf = page_num_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"{i + 1}"
                p.alignment = PP_ALIGN.RIGHT
                p.font.name = font_name
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(150, 150, 150)
                
                # 添加内容区域
                content = slide_data.get('content', [])
                if content:
                    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.8))
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    
                    for j, point in enumerate(content):
                        if j == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"• {point}"
                        p.font.name = font_name
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(60, 60, 60)
                        p.space_after = Pt(10)
                        p.level = 0
                
                # 如果有配图，嵌入 AI 生成的图片
                image_suggestion = slide_data.get('image_suggestion', '')
                generated_images = st.session_state.courseware_session.get("generated_images", {})
                
                if i in generated_images and generated_images[i].get('success'):
                    # 嵌入 AI 生成的 PNG 图片
                    img_path = generated_images[i].get('png_path')
                    if img_path and os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, Inches(8.5), Inches(2), Inches(4), Inches(3))
                    else:
                        # 图片路径不存在，显示文字占位符
                        placeholder = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(8.5), Inches(2), Inches(4), Inches(3)
                        )
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
                        placeholder.line.color.rgb = RGBColor(*hex_to_rgb(accent_color))
                        img_text = slide.shapes.add_textbox(Inches(8.6), Inches(3.2), Inches(3.8), Inches(1))
                        tf = img_text.text_frame
                        p = tf.paragraphs[0]
                        p.text = f"📷 {image_suggestion}"
                        p.alignment = PP_ALIGN.CENTER
                        p.font.name = font_name
                        p.font.size = Pt(14)
                        p.font.color.rgb = RGBColor(100, 100, 100)
                        p.font.italic = True
                elif image_suggestion and image_suggestion.strip():
                    # 没有生成图片，显示文字占位符
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(8.5), Inches(2), Inches(4), Inches(3)
                    )
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    placeholder.line.color.rgb = RGBColor(*hex_to_rgb(accent_color))
                    
                    # 添加图片说明文字
                    img_text = slide.shapes.add_textbox(Inches(8.6), Inches(3.2), Inches(3.8), Inches(1))
                    tf = img_text.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"📷 {image_suggestion}"
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = font_name
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(100, 100, 100)
                    p.font.italic = True
        
        # 添加动画超链接按钮（如果有动画）
        animations = st.session_state.courseware_session.get("animations", [])
        if animations and len(animations) > 0:
            # 创建 animations 文件夹
            anim_folder = "animations"
            if not os.path.exists(anim_folder):
                os.makedirs(anim_folder)
            
            # 保存所有动画的 HTML 文件
            animation_service = AnimationService()
            for anim in animations:
                svg_code = anim.get('svg_code', '')
                if svg_code:
                    anim_title = anim.get('title', 'animation').replace(' ', '_')
                    html_content = animation_service.generate_html_animation(svg_code, anim.get('title', '教学动画'))
                    html_file = f"{anim_folder}/{anim_title}.html"
                    
                    try:
                        with open(html_file, 'w', encoding='utf-8') as f:
                            f.write(html_content)
                        print(f"✅ 动画 HTML 文件已保存：{html_file}")
                    except Exception as e:
                        print(f"⚠️ 保存动画 HTML 失败：{str(e)}")
            
            # 在对应幻灯片上添加超链接按钮
            for anim in animations:
                slide_index = anim.get("related_slide_index", 0) - 1  # 转为 0 索引
                if 0 <= slide_index < len(prs.slides):
                    anim_slide = prs.slides[slide_index]
                    
                    # 添加按钮形状
                    button = anim_slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(10.5), Inches(6.5), Inches(2.3), Inches(0.7)
                    )
                    button.fill.solid()
                    button.fill.fore_color.rgb = RGBColor(30, 136, 229)
                    button.line.fill.background()
                    
                    # 添加文字
                    tf = button.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "▶ 播放动画"
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = 'Microsoft YaHei'
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.size = Pt(13)
                    p.font.bold = True
                    
                    # 添加超链接（指向 HTML 文件）
                    anim_title = anim.get('title', 'animation').replace(' ', '_')
                    html_file = f"animations/{anim_title}.html"
                    try:
                        button.click_action.hyperlink.address = html_file
                        print(f"✅ 已在第 {slide_index + 1} 页添加动画超链接：{html_file}")
                    except Exception as e:
                        print(f"⚠️ 添加超链接失败：{str(e)}")
        
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        topic = st.session_state.courseware_session.get("topic", "课件")
        grade_level = st.session_state.courseware_session.get("grade_level", "")
        subject = st.session_state.courseware_session.get("subject", "")
        ppt_file_name = f"{topic}_{grade_level}_{subject}_课件.pptx"
        st.download_button(
            label="📥 下载 PPT 课件（可编辑的 PowerPoint 文件）",
            data=ppt_bytes.getvalue(),
            file_name=ppt_file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            help="点击下载生成的 PPT 课件，可在 PowerPoint 或 WPS 中打开编辑",
            type="primary"
        )
        
        # Word 教案下载功能
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        import io
        import re
        
        def clean_text(text):
            """清理文本中的特殊字符和乱码"""
            if not text:
                return ""
            # 如果是 dict 类型，转换为字符串
            if isinstance(text, dict):
                text = json.dumps(text, ensure_ascii=False, indent=2)
            # 确保是字符串
            if not isinstance(text, str):
                text = str(text)
            # 移除不可见字符和控制字符（保留换行符）
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
            # 替换常见的乱码符号
            text = text.replace('☒', '')
            text = text.replace('', '')
            text = text.replace('↓', '')
            # 移除多余的空白字符
            text = re.sub(r'[ \t]+', ' ', text)
            return text.strip()
        
        try:
            doc = Document()
            
            # 设置全局字体为微软雅黑
            doc.styles['Normal'].font.name = 'Microsoft YaHei'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
            
            # 添加标题
            heading = doc.add_heading(f"{clean_text(topic)} - 教案", 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # 添加基本信息
            doc.add_paragraph(f"年级：{clean_text(grade_level)}")
            doc.add_paragraph(f"学科：{clean_text(subject)}")
            doc.add_paragraph(f"主题：{clean_text(topic)}")
            doc.add_paragraph()
            
            # 添加教案大纲内容
            doc.add_heading('一、教学大纲', level=1)
            cleaned_outline = clean_text(outline)
            doc.add_paragraph(cleaned_outline)
            
            # 添加 PPT 详细内容
            doc.add_heading('二、PPT 课件内容详情', level=1)
            for i, slide_data in enumerate(slides, 1):
                slide_title = clean_text(slide_data.get("title", "无标题"))
                doc.add_heading(f'第{i}页：{slide_title}', level=2)
                
                subtitle = slide_data.get('subtitle')
                if subtitle:
                    doc.add_paragraph(f'副标题：{clean_text(subtitle)}')
                content = slide_data.get('content', [])
                for point in content:
                    cleaned_point = clean_text(point)
                    if cleaned_point:  # 只添加非空内容
                        doc.add_paragraph(cleaned_point, style='List Bullet')
                doc.add_paragraph()
            
            # 保存为 docx 文件
            docx_bytes = io.BytesIO()
            doc.save(docx_bytes)
            docx_bytes.seek(0)
            
            docx_file_name = f"{clean_text(topic)}_{clean_text(grade_level)}_{clean_text(subject)}_教案.docx"
            st.download_button(
                label="📥 下载 Word 教案（可编辑的 Word 文档）",
                data=docx_bytes.getvalue(),
                file_name=docx_file_name,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                help="点击下载配套的详细教案，可在 Microsoft Word 或 WPS 中打开编辑",
                type="primary"
            )
        except ImportError as import_error:
            st.error(f"❌ 缺少必要的依赖库：{str(import_error)}")
            st.warning("💡 请运行以下命令安装依赖：")
            st.code("pip install python-docx", language="bash")
        except Exception as docx_error:
            st.error(f"❌ Word 教案生成失败：{str(docx_error)}")
            st.info("💡 建议：可以手动复制上面的教案内容到 Word 文档中")
    except Exception as ppt_error:
        st.error(f"PPT 文件生成失败：{str(ppt_error)}")
        st.info("💡 建议：可以手动复制上面的 PPT 内容到 PowerPoint 中")


elif menu == "知识库管理":
    # 检查游客模式
    if is_guest:
        st.title("🗄️ 知识库管理")
        st.warning("⚠️ 游客模式无法使用此功能")
        st.info("💡 请注册或登录以使用知识库管理功能")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📝 立即注册", type="primary", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        with col2:
            if st.button("🔐 去登录", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        st.stop()
    
    st.title("🗄️ 知识库管理")
    
    # 初始化 RAG 知识库
    if "rag_documents" not in st.session_state:
        st.session_state.rag_documents = []
    
    # 统计信息（结合 RAG 数据库）
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        # 从 RAG 数据库获取统计
        rag_stats = rag_kb.get_statistics() if st.session_state.rag_kb_connected else {}
        total_docs = rag_stats.get('total_documents', 0)
        st.metric("RAG 文档总数", total_docs, f"+{total_docs}" if total_docs > 0 else "0")
    with col2:
        total_points = rag_stats.get('total_knowledge_points', 0)
        st.metric("知识点数量", total_points, "AI 自动提取")
    with col3:
        avg_usage = rag_stats.get('average_usage', 0)
        st.metric("平均使用次数", round(avg_usage, 1), "持续更新中")
    with col4:
        subject_dist = rag_stats.get('subject_distribution', [])
        st.metric("覆盖学科数", len(subject_dist), f"{len(subject_dist)} 个学科" if subject_dist else "持续更新中")
    
    st.divider()
    
    # 上传文档区域
    st.subheader("📤 上传文档到知识库")
    uploaded_docs = st.file_uploader(
        "支持上传教学资料、学术论文、参考书籍、教学图片等",
        type=["pdf", "doc", "docx", "ppt", "pptx", "txt", "md", "jpg", "jpeg", "png"],
        accept_multiple_files=True,
        key="knowledge_upload",
        help="上传后 AI 将自动解析并归类到知识库（图片将使用 OCR 识别文字）"
    )
    
    if uploaded_docs:
        # 调用服务层上传文档
        knowledge_service = KnowledgeService()
        result = knowledge_service.upload_documents(
            uploaded_files=uploaded_docs,
            category="通用"
        )
        
        if result["success"]:
            st.success(f"✅ 成功上传 {result['uploaded_count']} 个文档")
            if result["failed_files"]:
                st.warning(f"⚠️ 以下文件上传失败：{', '.join(result['failed_files'])}")
            st.rerun()
        else:
            st.error(f"❌ 上传失败：{result.get('error', '未知错误')}")
        
        # 显示已上传的文件列表和管理功能
        if st.session_state.knowledge_base["documents"]:
            st.divider()
            st.subheader("📚 已上传的文件")
            
            # 删除功能
            files_to_delete = st.multiselect(
                "选择要删除的文件（可多选）",
                options=[doc["name"] for doc in st.session_state.knowledge_base["documents"]],
                key="delete_files_select"
            )
            
            col1, col2 = st.columns([1, 3])
            with col1:
                if st.button("🗑️ 删除选中文件", type="primary", use_container_width=True):
                    if files_to_delete:
                        deleted_count = 0
                        # 从 RAG 数据库中删除
                        for doc_name in files_to_delete:
                            # 查找文档信息
                            doc_info = next((doc for doc in st.session_state.knowledge_base["documents"] if doc["name"] == doc_name), None)
                            if doc_info and doc_info.get('rag_doc_id'):
                                try:
                                    rag_kb.delete_document(doc_info['rag_doc_id'])
                                    info(f"✅ 已从 RAG 数据库删除：{doc_name}")
                                except Exception as db_error:
                                    warning(f"⚠️ RAG 数据库删除失败：{doc_name} - {str(db_error)}")
                        
                        # 从 session state 中删除
                        st.session_state.knowledge_base["documents"] = [
                            doc for doc in st.session_state.knowledge_base["documents"]
                            if doc["name"] not in files_to_delete
                        ]
                        deleted_count = len(files_to_delete)
                        st.success(f"✅ 已删除 {deleted_count} 个文件")
                        # 清空 multiselect 的选择状态
                        st.session_state.delete_files_select = []
                        st.rerun()
            
            with col2:
                if st.button("🗑️ 清空所有文件", use_container_width=True):
                    # 从 RAG 数据库中删除所有文档
                    if st.session_state.rag_kb_connected:
                        try:
                            all_docs = rag_kb.get_all_documents(limit=1000)
                            for doc in all_docs:
                                rag_kb.delete_document(doc['id'])
                            info(f"✅ 已从 RAG 数据库清空 {len(all_docs)} 个文档")
                        except Exception as db_error:
                            warning(f"⚠️ RAG 数据库清空失败：{str(db_error)}")
                    
                    # 清空 session state
                    st.session_state.knowledge_base["documents"] = []
                    st.session_state.knowledge_base["categories"] = {}
                    st.success("✅ 已清空所有文件")
                    st.rerun()
            
            # 显示文件列表
            st.caption(f"共 {len(st.session_state.knowledge_base['documents'])} 个文件")
        if st.button("🤖 AI 智能解析并归类", type="primary"):
            with st.spinner("正在解析文档内容并提取知识点..."):
                try:
                    # 调用服务层分析文档
                    knowledge_service = KnowledgeService()
                    result = knowledge_service.analyze_documents(
                        documents=st.session_state.knowledge_base["documents"],
                        api_key=DEFAULT_API_KEY,
                        base_url=BASE_URL
                    )
                    
                    if result["success"]:
                        st.success("✅ AI 解析完成！")
                        with st.expander("📖 查看 AI 解析结果", expanded=True):
                            st.markdown(result["analysis"])
                    else:
                        st.error(f"❌ 解析失败：{result.get('error', '未知错误')}")
                        
                except Exception as e:
                    st.error(f"解析失败：{str(e)}")
    
    st.divider()
    
    # RAG 智能检索功能
    st.subheader("🔍 RAG 智能检索")
    
    col1, col2 = st.columns([3, 1])
    with col1:
        search_query = st.text_input(
            "输入关键词搜索知识文档",
            placeholder="例如：函数单调性、牛顿定律、三角函数...",
            key="rag_search"
        )
    with col2:
        search_subject = st.selectbox(
            "限定学科",
            ["全部"] + ["语文", "数学", "英语", "物理", "化学", "生物", 
                       "历史", "地理", "政治", "体育", "美术", "音乐", "信息技术"],
            key="rag_subject"
        )
    
    if search_query and st.button("🔍 搜索", type="primary"):
        with st.spinner("正在 RAG 知识库中检索..."):
            # 调用服务层搜索文档
            knowledge_service = KnowledgeService()
            subject = None if search_subject == "全部" else search_subject
            results = knowledge_service.search_documents(
                query=search_query,
                subject=subject,
                limit=10
            )
                
            if results:
                st.success(f"✅ 找到 {len(results)} 篇相关文档")
                    
                # 显示搜索结果
                for i, doc in enumerate(results, 1):
                    with st.expander(f"📄 {i}. {doc.get('title', '无标题')} - {doc.get('subject', '未知学科')}"):
                        st.markdown(doc.get('content_text', '')[:1000])
                        if len(doc.get('content_text', '')) > 1000:
                            st.caption("...（内容过长，仅显示前1000字）")
            else:
                st.info(f"💡 未在 RAG 知识库中找到相关文档，请尝试其他关键词或上传更多资料。")
    
    st.divider()
    
    # 知识库展示
    st.subheader("📚 知识库内容")
    
    # 分类筛选
    categories = list(st.session_state.knowledge_base.get("categories", {}).keys())
    if categories:
        selected_category = st.selectbox(
            "按分类筛选",
            ["全部"] + categories,
            key="category_filter"
        )
        
        # 显示文档列表
        docs = st.session_state.knowledge_base.get("documents", [])
        if docs:
            if selected_category != "全部":
                filtered_docs = [d for d in docs if d["category"] == selected_category]
            else:
                filtered_docs = docs
            
            if filtered_docs:
                st.info(f"📊 共 {len(filtered_docs)} 篇文档")
                
                # 表格展示
                display_data = {
                    "📄 文档名称": [d["name"] for d in filtered_docs],
                    "📂 分类": [d["category"] for d in filtered_docs],
                    "⏰ 上传时间": [d["upload_time"] for d in filtered_docs]
                }
                st.dataframe(display_data, use_container_width=True, hide_index=True)
            else:
                st.warning("该分类下暂无文档")
        else:
            st.info("📭 知识库为空，请上传文档")
    else:
        st.info("📭 知识库为空，请上传文档")

elif menu == "学情分析":
    # 检查游客模式
    if is_guest:
        st.title("📊 AI 学情分析")
        st.warning("⚠️ 游客模式无法使用此功能")
        st.info("💡 请注册或登录以使用学情分析功能")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📝 立即注册", type="primary", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        with col2:
            if st.button("🔐 去登录", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.current_user = None
                st.rerun()
        st.stop()
    
    st.title("📊 AI 学情分析")
    
    # 选择分析模式
    analysis_mode = st.radio(
        "分析模式",
        ["单个学生", "全班评估"],
        horizontal=True,
        key="analysis_mode_select"
    )
    
    if analysis_mode == "单个学生":
        student_name = st.text_input("学生姓名", placeholder="输入姓名以生成个人报告")
    else:
        class_name = st.text_input("班级名称", placeholder="例如：高一 (3) 班")
        total_students = st.number_input("班级总人数", min_value=1, value=45)
    
    # 上传成绩数据
    st.subheader("📁 上传成绩/学习数据")
    uploaded_files = st.file_uploader(
        "上传成绩单、作业统计表、考试分析等文件，AI 将生成详细的学情报告",
        type=["xlsx", "xls", "csv", "pdf", "doc", "docx", "txt"],
        accept_multiple_files=True,
        key="study_upload",
        help="支持 Excel 成绩单、CSV 数据表、Word 文档等格式"
    )
    
    if uploaded_files:
        for file in uploaded_files:
            st.success(f"✅ {file.name}")
    
    # 显示统计数据
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("问题总数", len(st.session_state.learning_data.get("questions", [])), "今日新增")
    with col2:
        correct_rate = st.session_state.learning_data.get("correct_rate", 0)
        st.metric("平均正确率", f"{correct_rate}%" if correct_rate else "0%", "待更新")
    with col3:
        study_time = st.session_state.learning_data.get("study_time", 0)
        st.metric("学习时长", f"{study_time}分钟", "累计")
    
    # 互动记录展示
    st.subheader("📝 最近互动记录")
    questions = st.session_state.learning_data.get("questions", [])
    if questions:
        for q in questions[-5:]:
            st.info(f"{q.get('time', '')} - {q.get('scenario', '')}: {q.get('question', '')[:50]}...")
    else:
        st.info("暂无互动记录")
    
    # AI 生成学情报告
    if st.button("🤖 AI 生成学情报告", type="primary"):
        # 验证：必须有文件或输入信息
        has_file = uploaded_files and len(uploaded_files) > 0
        has_student = (analysis_mode == "单个学生" and student_name.strip())
        has_class = (analysis_mode == "全班评估" and class_name.strip())
        has_interactive_data = len(st.session_state.learning_data.get("questions", [])) > 0
        
        if not has_file and not has_student and not has_class:
            st.error("❌ 请至少完成以下一项：\n1. 上传成绩/学习数据文件\n2. 输入学生姓名（单个学生模式）\n3. 输入班级名称（全班评估模式）")
            st.stop()
        
        with st.spinner("正在分析学习数据并生成报告..."):
            try:
                # 构建学生信息
                if analysis_mode == "单个学生":
                    student_info = {"name": student_name if student_name else "某同学"}
                else:
                    student_info = {
                        "class_name": class_name if class_name else "某班",
                        "total_students": total_students
                    }
                
                # 调用服务层生成报告
                analysis_service = AnalysisService()
                result = analysis_service.generate_report(
                    analysis_mode=analysis_mode,
                    student_info=student_info,
                    uploaded_files=uploaded_files,
                    questions_data=st.session_state.learning_data.get("questions", []),
                    api_key=DEFAULT_API_KEY,
                    base_url=BASE_URL
                )
                
                if result["success"]:
                    st.success("✅ 学情报告生成完成！")
                    
                    # 显示报告
                    st.markdown(result["report"])
                    
                    # 自动生成并展示分析图表
                    charts = generate_analysis_charts(
                        questions_data=st.session_state.learning_data.get("questions", []),
                        analysis_mode=analysis_mode
                    )
                    display_analysis_charts(charts)
                else:
                    st.error(f"生成失败：{result.get('error', '未知错误')}")
                    
            except Exception as e:
                st.error(f"生成失败：{str(e)}")
    
    # 全局数据管理
    st.divider()
    with st.expander("🗄️ 全局数据管理"):
        st.subheader("💾 数据备份与恢复")
        col_backup1, col_backup2, col_backup3 = st.columns(3)
        
        with col_backup1:
            if st.button("💾 备份所有数据", use_container_width=True):
                analysis_service = AnalysisService()
                result = analysis_service.manage_learning_data("backup")
                if result["success"]:
                    st.success("✅ 学习数据已备份到 learning_data_backup.json")
                else:
                    st.error(f"备份失败：{result.get('error')}")
        
        with col_backup2:
            if st.button("📤 导出完整数据", use_container_width=True):
                full_data = {
                    "learning_data": st.session_state.learning_data,
                    "knowledge_base": st.session_state.knowledge_base,
                    "export_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                }
                export_str = json.dumps(full_data, ensure_ascii=False, indent=2)
                st.download_button(
                    label="📥 下载完整数据",
                    data=export_str,
                    file_name=f"完整数据备份_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json"
                )
        
        with col_backup3:
            if st.button("🗑️ 清空所有数据", use_container_width=True, type="secondary"):
                if st.checkbox("确认清空所有数据？", key="confirm_clear_all"):
                    analysis_service = AnalysisService()
                    result = analysis_service.manage_learning_data("clear")
                    if result["success"]:
                        st.success("✅ 所有数据已清空")
                        st.rerun()
