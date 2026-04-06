"""
文档解析与 JSON 格式转换模块
支持多种文件格式转换为统一的 JSON 格式
"""

import json
from datetime import datetime
import os


class DocumentParser:
    """文档解析器 - 将各种文件格式转换为统一的 JSON 格式"""
    
    @staticmethod
    def parse_to_json(file, subject="通用", uploaded_by="teacher"):
        """
        将上传的文件解析为统一的 JSON 格式
        
        参数：
        - file: Streamlit UploadedFile 对象
        - subject: 学科分类（语文、数学、英语等）
        - uploaded_by: 上传者
        
        返回：
        - dict: 标准化的 JSON 格式数据
        """
        try:
            # 获取文件基本信息
            file_name = file.name
            file_ext = file_name.split('.')[-1].lower()
            file_size = file.size
            
            # 读取文件内容
            content_text = ""
            if file_ext in ['txt', 'md']:
                content_text = file.read().decode('utf-8')
            elif file_ext in ['doc', 'docx']:
                content_text = DocumentParser._read_docx(file)
            elif file_ext == 'pdf':
                content_text = DocumentParser._read_pdf(file)
            elif file_ext in ['ppt', 'pptx']:
                content_text = DocumentParser._read_pptx(file)
            
            # 重置文件指针，以便后续下载
            file.seek(0)
            
            # 构建统一的 JSON 结构
            document_data = {
                "metadata": {
                    "title": file_name,
                    "subject": subject,
                    "file_type": file_ext,
                    "file_size": file_size,
                    "uploaded_by": uploaded_by,
                    "upload_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "version": "1.0"
                },
                "content": {
                    "raw_text": content_text[:50000],  # 限制长度
                    "text_length": len(content_text),
                    "paragraphs": DocumentParser._split_paragraphs(content_text),
                    "word_count": len(content_text.split())
                },
                "analysis": {
                    "knowledge_points": [],  # AI 提取的知识点
                    "summary": "",  # AI 生成的摘要
                    "difficulty_level": "中等",  # 难度等级
                    "tags": []  # 标签
                },
                "storage": {
                    "original_filename": file_name,
                    "stored_format": "json",
                    "compression": False
                }
            }
            
            return document_data
            
        except Exception as e:
            print(f"❌ 解析文件失败：{str(e)}")
            return None
    
    @staticmethod
    def _read_docx(file):
        """读取 Word 文档"""
        try:
            from docx import Document
            import io
            
            doc = Document(io.BytesIO(file.read()))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n".join(paragraphs)
        except ImportError:
            return "[需要安装 python-docx 库]"
        except Exception as e:
            return f"[Word 文档读取失败：{str(e)}]"
    
    @staticmethod
    def _read_pdf(file):
        """读取 PDF 文档"""
        try:
            from PyPDF2 import PdfReader
            import io
            
            pdf_reader = PdfReader(io.BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            return text
        except ImportError:
            return "[需要安装 PyPDF2 库]"
        except Exception as e:
            return f"[PDF 文档读取失败：{str(e)}]"
    
    @staticmethod
    def _read_pptx(file):
        """读取 PPT 文档"""
        try:
            from pptx import Presentation
            import io
            
            prs = Presentation(io.BytesIO(file.read()))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slides_text.append(text)
            return "\n\n".join(slides_text)
        except ImportError:
            return "[需要安装 python-pptx 库]"
        except Exception as e:
            return f"[PPT 文档读取失败：{str(e)}]"
    
    @staticmethod
    def _split_paragraphs(text, max_length=500):
        """将文本分割为段落"""
        if not text:
            return []
        
        # 按换行符分割
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        # 如果段落太长，进一步分割
        result = []
        for para in paragraphs:
            if len(para) <= max_length:
                result.append(para)
            else:
                # 按句子分割
                sentences = para.split('。')
                current = ""
                for sentence in sentences:
                    if len(current + sentence) <= max_length:
                        current += sentence + "。"
                    else:
                        if current:
                            result.append(current.strip())
                        current = sentence + "。"
                if current:
                    result.append(current.strip())
        
        return result[:100]  # 最多保留 100 个段落
    
    @staticmethod
    def to_json_string(document_data, indent=2):
        """将文档数据转换为 JSON 字符串"""
        try:
            return json.dumps(document_data, ensure_ascii=False, indent=indent)
        except Exception as e:
            print(f"❌ JSON 转换失败：{str(e)}")
            return "{}"
    
    @staticmethod
    def save_to_file(document_data, output_dir="uploads/json_docs"):
        """保存 JSON 文档到文件"""
        try:
            # 创建输出目录
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # 生成文件名
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            title = document_data["metadata"]["title"].replace(".", "_")
            filename = f"{timestamp}_{title}.json"
            filepath = os.path.join(output_dir, filename)
            
            # 写入 JSON 文件
            json_str = DocumentParser.to_json_string(document_data)
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(json_str)
            
            return filepath
            
        except Exception as e:
            print(f"❌ 保存 JSON 文件失败：{str(e)}")
            return None
    
    @staticmethod
    def load_from_file(filepath):
        """从 JSON 文件加载文档数据"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return data
        except Exception as e:
            print(f"❌ 加载 JSON 文件失败：{str(e)}")
            return None


# 创建全局解析器实例
doc_parser = DocumentParser()
